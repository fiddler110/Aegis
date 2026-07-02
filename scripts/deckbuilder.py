#!/usr/bin/env python3
"""
deckbuilder.py — Convert YAML or Markdown to polished PowerPoint presentations.

Amalgamates patterns from:
  addq3pislidetoexistingdeck.py, decisiondependencydiagramdeck.py,
  enhancedvisualexecdeckgen.py, fixedminimalroadmap.py,
  identity-correctedexecdeck.py, recommendedfuturetemplatescript.py,
  simpleexecroadmapdeckgen.py, simplifiedkeydecisionsandriskslide.py

Usage:
    python deckbuilder.py deck.yaml
    python deckbuilder.py deck.md -o output.pptx
    python deckbuilder.py deck.yaml --theme corporate_blue
    python deckbuilder.py --list-themes

Requires:
    pip install python-pptx pyyaml
"""

import argparse
import re
import sys
import textwrap
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ===========================================================================
# BUILT-IN COLOUR THEMES
# ===========================================================================

BUILT_IN_THEMES: dict[str, dict] = {
     "Corp_red": {
        "primary":    "#ED1C24",   # brand red   — primary accent
        "dark":       "#231F20",   # near-black body text
        "navy":       "#243558",   # deepened blue for dark slide bgs
        "blue":       "#3A517C",   # dark blue   — card / badge accent
        "teal":       "#009863",   # green       — positive / delivered
        "green":      "#009863",   # green
        "orange":     "#D47F4C",   # orange      — caution / medium risk
        "light_gray": "#F5F5F5",   # lightgrey   — card backgrounds
        "mid_gray":   "#646464",   # mid grey    — footer / subtitle text
        "white":      "#FFFFFF",
    },
    "corporate_red": {
        "primary":    "#CE1126",
        "dark":       "#22304C",
        "navy":       "#19375A",
        "blue":       "#2980B9",
        "teal":       "#008080",
        "green":      "#27AE60",
        "orange":     "#F39C12",
        "light_gray": "#F2F4F7",
        "mid_gray":   "#6C7578",
        "white":      "#FFFFFF",
    },
    "corporate_blue": {
        "primary":    "#1A5276",
        "dark":       "#1A2A3A",
        "navy":       "#0D2137",
        "blue":       "#2980B9",
        "teal":       "#1ABC9C",
        "green":      "#27AE60",
        "orange":     "#E67E22",
        "light_gray": "#EBF3FB",
        "mid_gray":   "#7F8C8D",
        "white":      "#FFFFFF",
    },
    "corporate_green": {
        "primary":    "#1E8449",
        "dark":       "#1A2A2A",
        "navy":       "#0D1F1F",
        "blue":       "#2E86C1",
        "teal":       "#148F77",
        "green":      "#27AE60",
        "orange":     "#D4AC0D",
        "light_gray": "#EAFAF1",
        "mid_gray":   "#7F8C8D",
        "white":      "#FFFFFF",
    },
    "modern_purple": {
        "primary":    "#8E44AD",
        "dark":       "#2C3E50",
        "navy":       "#1A252F",
        "blue":       "#2980B9",
        "teal":       "#16A085",
        "green":      "#27AE60",
        "orange":     "#E67E22",
        "light_gray": "#F4ECF7",
        "mid_gray":   "#95A5A6",
        "white":      "#FFFFFF",
    },
    "minimal_slate": {
        "primary":    "#2C3E50",
        "dark":       "#1A1A2E",
        "navy":       "#16213E",
        "blue":       "#4A90D9",
        "teal":       "#45B8AC",
        "green":      "#5BAD6F",
        "orange":     "#E8A838",
        "light_gray": "#F0F2F5",
        "mid_gray":   "#8E9BB3",
        "white":      "#FFFFFF",
    },
    "executive_dark": {
        "primary":    "#E74C3C",
        "dark":       "#ECF0F1",
        "navy":       "#1C2833",
        "blue":       "#3498DB",
        "teal":       "#1ABC9C",
        "green":      "#2ECC71",
        "orange":     "#F39C12",
        "light_gray": "#2C3E50",
        "mid_gray":   "#95A5A6",
        "white":      "#FFFFFF",
    },
}


# ===========================================================================
# THEME
# ===========================================================================

def _hex_to_rgb(hex_val: str) -> RGBColor:
    h = hex_val.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class Theme:
    """
    Resolved colour palette.

    Can be initialised with:
      - A string naming a built-in theme  ("corporate_red")
      - A dict with optional "base" key + per-colour overrides
        {"base": "corporate_blue", "primary": "#FF6600"}
    """

    def __init__(self, config):
        if isinstance(config, str):
            palette = dict(BUILT_IN_THEMES.get(config, BUILT_IN_THEMES["corporate_red"]))
        elif isinstance(config, dict):
            base_name = config.get("base", "corporate_red")
            palette = dict(BUILT_IN_THEMES.get(base_name, BUILT_IN_THEMES["corporate_red"]))
            palette.update({k: v for k, v in config.items() if k != "base"})
        else:
            palette = dict(BUILT_IN_THEMES["corporate_red"])

        for name, hex_val in palette.items():
            setattr(self, name, _hex_to_rgb(hex_val))

    def color(self, name_or_hex: str) -> RGBColor:
        """
        Resolve a colour reference.
          - Named palette key  → e.g. "primary", "teal", "orange"
          - Hex string         → e.g. "#FF5500"
          - Fallback           → theme primary
        """
        if not name_or_hex:
            return self.primary
        if str(name_or_hex).startswith("#"):
            return _hex_to_rgb(name_or_hex)
        return getattr(self, name_or_hex, self.primary)


# ===========================================================================
# RENDERING ENGINE
# ===========================================================================

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER_Y = Inches(7.05)


class DeckRenderer:
    """
    Converts structured slide definitions into a python-pptx Presentation.

    Supported slide types (see _RENDERERS for all aliases):
      title, bullets, cards, decision_risk, dependency,
      timeline, quote, metrics, divider, table
    """

    def __init__(self, theme: Theme, footer: str = "", *, logo_path: str = None,
                 text_offset: int = 0):
        self.t = theme
        self.footer = footer
        self.logo_path = logo_path
        self.text_offset = text_offset
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def _sz(self, base: float) -> Pt:
        """Return Pt(base + text_offset) — applies the global text-size adjustment."""
        return Pt(base + self.text_offset)

    # -----------------------------------------------------------------------
    # LOW-LEVEL DRAWING PRIMITIVES
    # -----------------------------------------------------------------------

    def _blank(self):
        """Return a completely blank slide (layout 6 = Blank)."""
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide, *, bg_color=None):
        """
        White background + thin primary top bar + optional footer text.
        bg_color overrides the background fill (for dark slides, etc.).
        """
        t = self.t
        fill_color = t.color(bg_color) if bg_color else t.white

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = fill_color
        bg.line.fill.background()

        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
        top.fill.solid()
        top.fill.fore_color.rgb = t.primary
        top.line.fill.background()

        if self.footer:
            box = slide.shapes.add_textbox(
                Inches(0.35), FOOTER_Y, Inches(12.6), Inches(0.25)
            )
            p = box.text_frame.paragraphs[0]
            p.text = self.footer
            p.font.size = self._sz(8)
            p.font.color.rgb = t.mid_gray
            p.alignment = PP_ALIGN.RIGHT

        if self.logo_path:
            try:
                # Logo sits just below the primary top bar, flush to the right
                # edge with a small padding.  Aspect ratio 1.145 → at height
                # 0.60" the rendered width ≈ 0.69"; left is calculated so the
                # right edge lands at SLIDE_W − 0.15".
                _LOGO_H = Inches(1.00)
                _LOGO_W = Inches(1.00 * 1.145)   # ~1.15" — preserves aspect ratio
                slide.shapes.add_picture(
                    self.logo_path,
                    SLIDE_W - _LOGO_W - Inches(0.15),  # right-aligned, slight padding
                    Inches(0.14),                       # just below the top bar
                    height=_LOGO_H,
                )
            except Exception:
                pass  # logo file missing or unreadable — silently skip

    def _heading(self, slide, title: str, subtitle: str = None):
        """Large bold title + optional italic subtitle line."""
        t = self.t
        tb = slide.shapes.add_textbox(
            Inches(0.35), Inches(0.22), Inches(12.6), Inches(0.6)
        )
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = self._sz(26)
        p.font.bold = True
        p.font.color.rgb = t.dark

        if subtitle:
            st = slide.shapes.add_textbox(
                Inches(0.37), Inches(0.85), Inches(12.3), Inches(0.32)
            )
            p2 = st.text_frame.paragraphs[0]
            p2.text = subtitle
            p2.font.size = self._sz(11)
            p2.font.color.rgb = t.mid_gray

    def _textbox(
        self, slide, x, y, w, h, text, *,
        size=11, color=None, bold=False, align=PP_ALIGN.LEFT, wrap=True
    ):
        """Add a plain text box at the given (x, y, w, h) position in inches."""
        color = color or self.t.dark
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = self._sz(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return tb

    def _badge(
        self, slide, x, y, w, h, text, fill, *,
        font_color=None, size=10, bold=True
    ):
        """
        Solid rounded-rectangle label with centred, vertically-middle text.
        Used for column headers, RAG buckets, and decision boxes.
        """
        font_color = font_color or self.t.white
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = fill

        tf = shp.text_frame
        tf.clear()
        tf.margin_left = Inches(0.07)
        tf.margin_right = Inches(0.07)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = self._sz(size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        return shp

    def _card(
        self, slide, x, y, w, h, title, bullets=None, *,
        fill=None, accent=None, title_color=None,
        title_size=15, bullet_size=11.5
    ):
        """
        Content card: rounded rect with a coloured left-edge accent bar,
        a bold title, and an optional bullet list.
        This is the core visual building block, lifted from the exec deck scripts.
        """
        t = self.t
        fill = fill or t.white
        accent = accent or t.primary
        title_color = title_color or t.dark

        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = RGBColor(220, 224, 230)

        # Coloured left edge bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        tf = shp.text_frame
        tf.clear()
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.10)

        p = tf.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = self._sz(title_size)
        p.font.color.rgb = title_color
        p.space_after = Pt(6)

        for b in (bullets or []):
            pp = tf.add_paragraph()
            pp.text = "• " + b
            pp.font.size = self._sz(bullet_size)
            pp.font.color.rgb = t.dark

        return shp

    # -----------------------------------------------------------------------
    # SLIDE TYPE RENDERERS
    # -----------------------------------------------------------------------

    def _render_title(self, sd: dict):
        """
        Title / cover slide.
        White background, primary-colour top bar + bottom bar,
        primary-colour title text, dark subtitle text.
        """
        t = self.t
        slide = self._blank()

        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = t.white
        bg.line.fill.background()

        # Top primary bar (thicker than content slides)
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = t.primary
        top.line.fill.background()

        # Bottom primary bar
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25)
        )
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = t.primary
        bottom.line.fill.background()

        tb = slide.shapes.add_textbox(
            Inches(0.65), Inches(2.1), Inches(12.0), Inches(2.0)
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd.get("title", "")
        p.font.size = self._sz(42)
        p.font.bold = True
        p.font.color.rgb = t.primary

        if sd.get("subtitle"):
            st = slide.shapes.add_textbox(
                Inches(0.65), Inches(4.3), Inches(12.0), Inches(0.65)
            )
            p2 = st.text_frame.paragraphs[0]
            p2.text = sd["subtitle"]
            p2.font.size = self._sz(20)
            p2.font.color.rgb = t.dark

        if sd.get("context"):
            cx = slide.shapes.add_textbox(
                Inches(0.65), Inches(5.1), Inches(12.0), Inches(0.35)
            )
            p3 = cx.text_frame.paragraphs[0]
            p3.text = sd["context"]
            p3.font.size = self._sz(13)
            p3.font.color.rgb = t.mid_gray

    def _render_bullets(self, sd: dict):
        """
        Standard title + indented bullet list.
        Supports nesting via leading spaces ("  sub item") or
        explicit level dicts: {text: "...", level: 1}
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        y_start = 1.38 if not sd.get("subtitle") else 1.48

        if sd.get("intro"):
            self._textbox(
                slide, 0.55, y_start, 12.2, 0.48, sd["intro"],
                size=14, color=t.dark, bold=True, align=PP_ALIGN.CENTER
            )
            y_start += 0.60

        items = sd.get("bullets", [])
        if not items:
            return

        tb = slide.shapes.add_textbox(
            Inches(0.55), Inches(y_start), Inches(12.2), Inches(6.8 - y_start)
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True

        for idx, item in enumerate(items):
            # Detect depth from leading spaces or explicit dict
            if isinstance(item, dict):
                depth, text = item.get("level", 0), item.get("text", "")
            else:
                text = item
                if item.startswith("    "):
                    depth, text = 2, item.strip()
                elif item.startswith("  "):
                    depth, text = 1, item.strip()
                else:
                    depth = 0

            chars = ["•", "◦", "▪"][min(depth, 2)]
            indent = "    " * depth

            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = f"{indent}{chars} {text}"
            p.font.size = self._sz(max(15 - depth * 1.5, 11))
            p.font.color.rgb = t.dark
            p.space_before = Pt(5 if depth == 0 else 2)

    def _render_cards(self, sd: dict):
        """
        N-column card layout (2–4 columns, typically 3).
        Each column is a _card() with title + bullets.

        Optional extensions:
          header_card:  dict with title, bullets, color — full-width card above columns.
          flow_boxes:   list of {title, color} dicts — horizontal arrow-connected boxes
                        rendered below the column cards.
          takeaway:     string — full-width summary card below everything.
          note:         string — plain italic note text at the very bottom.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        y_cursor = 1.38
        margin, gap = 0.55, 0.22

        # ── optional intro text ────────────────────────────────────────────
        if sd.get("intro"):
            self._textbox(
                slide, 0.55, y_cursor, 12.2, 0.52, sd["intro"],
                size=14, color=t.dark, bold=True, align=PP_ALIGN.CENTER
            )
            y_cursor += 0.60

        # ── optional full-width header card ───────────────────────────────
        hc = sd.get("header_card")
        if hc:
            hc_fill = t.color(hc.get("color", "navy"))
            hc_accent = t.color(hc.get("accent", hc.get("color", "navy")))
            hc_h = hc.get("height", 0.80)
            self._card(
                slide, margin, y_cursor, 13.333 - 2 * margin, hc_h,
                hc.get("title", ""), hc.get("bullets", []),
                fill=hc_fill, accent=hc_accent,
                title_color=t.white,
                title_size=hc.get("title_size", 15),
                bullet_size=hc.get("bullet_size", 11.5),
            )
            y_cursor += hc_h + 0.18

        # ── main column cards ─────────────────────────────────────────────
        cols = sd.get("columns", [])
        default_accents = ["primary", "blue", "green", "teal", "orange", "navy"]

        if cols:
            n = len(cols)
            card_w = (13.333 - 2 * margin - gap * (n - 1)) / n

            # Reserve space for optional rows below
            has_flow   = bool(sd.get("flow_boxes"))
            has_take   = bool(sd.get("takeaway"))
            has_note   = bool(sd.get("note"))
            bottom_reserve = (
                (0.80 if has_flow else 0)
                + (0.82 if has_take else 0)
                + (0.28 if has_note else 0)
            )
            max_bottom = 7.05 - bottom_reserve
            card_h = sd.get("card_height", min(max_bottom - y_cursor - 0.15, 5.2))

            for i, col in enumerate(cols):
                x = margin + i * (card_w + gap)
                accent = t.color(col.get("accent", default_accents[i % len(default_accents)]))
                fill = t.color(col.get("fill", "white"))
                self._card(
                    slide, x, y_cursor, card_w, card_h,
                    col.get("title", ""), col.get("bullets", []),
                    fill=fill, accent=accent,
                    title_size=col.get("title_size", 15),
                    bullet_size=col.get("bullet_size", 11.5),
                )

            y_cursor += card_h + 0.18

        # ── optional flow boxes (arrow-connected process flow) ─────────────
        flow_boxes = sd.get("flow_boxes", [])
        if flow_boxes:
            n_fb = len(flow_boxes)
            arrow_w = 0.38
            total_arrow = arrow_w * (n_fb - 1)
            fb_w = (13.333 - 2 * margin - gap * (n_fb - 1) - total_arrow) / n_fb
            fb_h = sd.get("flow_height", 0.60)
            fb_x = margin
            for j, fb in enumerate(flow_boxes):
                fill = t.color(fb.get("color", default_accents[j % len(default_accents)]))
                self._badge(slide, fb_x, y_cursor, fb_w, fb_h, fb.get("title", ""), fill, size=10)
                fb_x += fb_w
                if j < n_fb - 1:
                    mid_y = y_cursor + fb_h / 2
                    self._textbox(
                        slide, fb_x + gap / 2, mid_y - 0.18, arrow_w - gap, 0.36, "→",
                        size=18, color=t.mid_gray, bold=True, align=PP_ALIGN.CENTER
                    )
                    fb_x += gap + arrow_w
            y_cursor += fb_h + 0.18

        # ── optional full-width takeaway card ──────────────────────────────
        if sd.get("takeaway"):
            tw_title = sd.get("takeaway_title", "Key takeaway")
            tw_text  = sd["takeaway"]
            self._card(
                slide, 0.85, y_cursor, 11.65, 0.72,
                tw_title, [tw_text],
                fill=t.light_gray, accent=t.primary,
                title_size=13, bullet_size=12,
            )
            y_cursor += 0.82

        # ── optional plain note text ───────────────────────────────────────
        if sd.get("note"):
            self._textbox(
                slide, 0.55, y_cursor, 12.2, 0.30, sd["note"],
                size=9, color=t.mid_gray, align=PP_ALIGN.LEFT
            )

    def _render_decision_risk(self, sd: dict):
        """
        Executive decisions + RAG risk grid.
        Decision boxes render across the top; colour-coded risk buckets below.
        Implements the pattern from decisiondependencydiagramdeck.py and
        simplifiedkeydecisionsandriskslide.py.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        decisions = sd.get("decisions", [])
        dec_bottom = 1.38

        if decisions:
            n = len(decisions)
            margin, gap = 0.55, 0.25
            dw = (13.333 - 2 * margin - gap * (n - 1)) / n
            dh = 1.15
            x = margin
            for dec in decisions:
                fill = t.color(dec.get("color", "primary"))
                text = dec.get("title", "")
                if dec.get("body"):
                    text += "\n" + dec["body"]
                self._badge(slide, x, 1.38, dw, dh, text, fill, size=7.5)
                x += dw + gap
            dec_bottom = 1.38 + dh + 0.22

        risks = sd.get("risks", {})
        risk_levels = [
            ("HIGH",   t.color(risks.get("high_color",   "primary")), risks.get("high",   [])),
            ("MEDIUM", t.color(risks.get("medium_color", "orange")),  risks.get("medium", [])),
            ("LOW",    t.color(risks.get("low_color",    "green")),   risks.get("low",    [])),
        ]

        bucket_h = max(6.8 - dec_bottom - 0.05, 0.5)
        margin, gap = 0.55, 0.25
        bw = (13.333 - 2 * margin - gap * 2) / 3
        x = margin

        for level, fill, items in risk_levels:
            if items:
                text = level + "\n" + "\n".join(items)
            else:
                text = level + "\n(none)"
            self._badge(slide, x, dec_bottom, bw, bucket_h, text, fill, size=9)
            x += bw + gap

    def _render_card_grid(self, sd: dict):
        """
        Multi-row card grid with two rendering modes.

        Mode 1 — simple grid (``rows`` key):
          Renders a 2-D array of content cards, e.g. 2 rows × 3 columns.
          Each row entry is a dict with an optional ``label`` / ``accent`` and
          a ``columns`` list of card dicts (title, accent, fill, bullets).

          Optional extras (all modes):
            header_badge: {title, color} — centred badge above all cards.
            badges:       list of short label strings shown as a badge row
                          below the cards (e.g. outcome summary pills).
            takeaway:     string — full-width summary card at the bottom.
            note:         string — plain italic note text at the very bottom.

        Mode 2 — two labelled side-by-side sections (``sections`` key):
          Each section has a ``label``, ``color`` (header fill), and its own
          ``rows`` list (same inner structure as Mode 1).  The two sections
          split the slide width evenly.  Used for "2027 / 2028" layouts.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        margin, gap = 0.55, 0.22
        y_cursor = 1.38
        default_accents = ["primary", "blue", "green", "teal", "orange", "navy"]

        # ── optional intro text ────────────────────────────────────────────
        if sd.get("intro"):
            self._textbox(
                slide, margin, y_cursor, 12.2, 0.40, sd["intro"],
                size=12.5, color=t.dark, bold=True, align=PP_ALIGN.CENTER
            )
            y_cursor += 0.50

        # ── optional centred header badge ──────────────────────────────────
        hb = sd.get("header_badge")
        if hb:
            hb_fill = t.color(hb.get("color", "primary"))
            hb_w = hb.get("width", 5.25)
            hb_x = (13.333 - hb_w) / 2
            self._badge(slide, hb_x, y_cursor, hb_w, 0.48,
                        hb.get("title", ""), hb_fill, size=14)
            y_cursor += 0.62

        # ── reserve space for bottom elements ─────────────────────────────
        has_badges   = bool(sd.get("badges"))
        has_takeaway = bool(sd.get("takeaway"))
        has_note     = bool(sd.get("note"))
        bottom_reserve = (
            (0.62 if has_badges   else 0)
            + (0.82 if has_takeaway else 0)
            + (0.32 if has_note     else 0)
        )

        # ── MODE 2 — two labelled sections ────────────────────────────────
        sections = sd.get("sections")
        if sections:
            n_sections = len(sections)
            section_gap = 0.35
            sw = (13.333 - 2 * margin - section_gap * (n_sections - 1)) / n_sections

            # Figure out the tallest section to unify card heights
            max_rows = max(len(sec.get("rows", [])) for sec in sections)
            max_cols = max(
                len(row.get("columns", []))
                for sec in sections
                for row in sec.get("rows", [])
            ) if any(sec.get("rows") for sec in sections) else 1

            label_h    = 0.48
            available  = 7.05 - bottom_reserve - y_cursor - label_h - 0.10
            row_gap    = 0.14
            card_h     = (available - row_gap * (max_rows - 1)) / max(max_rows, 1)

            for si, sec in enumerate(sections):
                sx = margin + si * (sw + section_gap)
                acc = t.color(sec.get("color", default_accents[si % len(default_accents)]))

                # Section header label
                self._badge(slide, sx, y_cursor, sw, label_h,
                            sec.get("label", f"Section {si+1}"), acc, size=13)

                row_y = y_cursor + label_h + 0.10
                sec_rows = sec.get("rows", [])
                n_cols_sec = max((len(r.get("columns", [])) for r in sec_rows), default=1)
                cw = (sw - gap * (n_cols_sec - 1)) / max(n_cols_sec, 1)

                for row in sec_rows:
                    for ci, col in enumerate(row.get("columns", [])):
                        cx = sx + ci * (cw + gap)
                        ca = t.color(col.get("accent",
                                             default_accents[ci % len(default_accents)]))
                        cf = t.color(col.get("fill", "white"))
                        self._card(
                            slide, cx, row_y, cw, card_h,
                            col.get("title", ""), col.get("bullets", []),
                            fill=cf, accent=ca,
                            title_size=col.get("title_size", 13.5),
                            bullet_size=col.get("bullet_size", 10.2),
                        )
                    row_y += card_h + row_gap

            y_cursor += label_h + 0.10 + max_rows * (card_h + row_gap)

        # ── MODE 1 — simple row grid ───────────────────────────────────────
        else:
            rows = sd.get("rows", [])
            if rows:
                # Determine uniform card geometry
                max_cols  = max(len(r.get("columns", [])) for r in rows)
                n_rows    = len(rows)
                cw        = (13.333 - 2 * margin - gap * (max_cols - 1)) / max(max_cols, 1)
                available = 7.05 - bottom_reserve - y_cursor
                row_gap   = 0.18
                card_h    = sd.get("card_height",
                                   (available - row_gap * (n_rows - 1)) / max(n_rows, 1))
                card_h    = min(card_h, 2.4)

                for row in rows:
                    cols = row.get("columns", [])
                    row_accent = t.color(row.get("accent", "primary"))
                    for ci, col in enumerate(cols):
                        cx = margin + ci * (cw + gap)
                        ca = t.color(col.get("accent",
                                             default_accents[ci % len(default_accents)]))
                        cf = t.color(col.get("fill", "white"))
                        self._card(
                            slide, cx, y_cursor, cw, card_h,
                            col.get("title", ""), col.get("bullets", []),
                            fill=cf, accent=ca,
                            title_size=col.get("title_size", 14.5),
                            bullet_size=col.get("bullet_size", 10.8),
                        )
                    y_cursor += card_h + row_gap

        # ── optional badge row (outcome pills) ────────────────────────────
        badges = sd.get("badges", [])
        if badges:
            n_b = len(badges)
            bw  = (13.333 - 2 * margin - gap * (n_b - 1)) / max(n_b, 1)
            bx  = margin
            for badge_text in badges:
                shp = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(bx), Inches(y_cursor), Inches(bw), Inches(0.46)
                )
                shp.fill.solid()
                shp.fill.fore_color.rgb = t.light_gray
                shp.line.color.rgb = RGBColor(220, 224, 230)
                tf = shp.text_frame
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = badge_text
                p.alignment = PP_ALIGN.CENTER
                p.font.size = self._sz(11)
                p.font.color.rgb = t.dark
                bx += bw + gap
            y_cursor += 0.56

        # ── optional full-width takeaway card ──────────────────────────────
        if sd.get("takeaway"):
            tw_title = sd.get("takeaway_title", "AVP takeaway")
            self._card(
                slide, margin, y_cursor, 13.333 - 2 * margin, 0.52,
                tw_title, [sd["takeaway"]],
                fill=t.light_gray, accent=t.primary,
                title_size=13, bullet_size=12,
            )
            y_cursor += 0.62

        # ── optional plain note text ───────────────────────────────────────
        if sd.get("note"):
            self._textbox(
                slide, margin, y_cursor, 12.2, 0.40, sd["note"],
                size=12, color=t.dark, bold=True, align=PP_ALIGN.CENTER
            )

    def _render_decisions_heatmap(self, sd: dict):
        """
        Side-by-side decisions panel + colour-coded risk heatmap table.

        Left column  — ``decisions`` dict:
          title, color (fill), items (list of bullet strings).
          Optional second box below: ``risk_themes`` dict with same shape.

        Right column — ``heatmap`` dict:
          title, headers (list), col_widths (optional list of inch widths),
          rows (list of {label, values} where each value is "High"/"Medium"/"Low"),
          level_colors (optional dict mapping level→palette key).
          Optional ``takeaway`` dict below: {title, text}.

        High / Medium / Low cells default to primary / orange / green.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        # ── column geometry ────────────────────────────────────────────────
        left_x, left_w  = 0.45, 4.10
        right_x         = 5.00
        right_w         = 13.333 - right_x - 0.45
        y_top           = 1.20

        # ── LEFT: decisions box ────────────────────────────────────────────
        dec   = sd.get("decisions", {})
        dec_h = dec.get("height", 2.15)
        if dec:
            dec_fill = t.color(dec.get("color", "primary"))
            bullets  = dec.get("items", [])
            self._card(
                slide, left_x, y_top, left_w, dec_h,
                dec.get("title", "Strategic decisions required"),
                bullets,
                fill=t.white, accent=dec_fill,
                title_size=dec.get("title_size", 14.5),
                bullet_size=dec.get("bullet_size", 10.5),
            )

        # ── RIGHT: heatmap ─────────────────────────────────────────────────
        hm = sd.get("heatmap", {})
        if hm:
            self._textbox(
                slide, right_x, y_top - 0.02, right_w, 0.30,
                hm.get("title", "Risk Heatmap Summary"),
                size=13, color=t.dark, bold=True
            )
            hdr_y    = y_top + 0.32
            headers  = hm.get("headers", [])
            hm_rows  = hm.get("rows", [])
            n_cols   = len(headers)
            # Default col widths: first col wider, rest equal
            default_cw = [2.45] + [round(right_w - 2.45) / max(n_cols - 1, 1)] * (n_cols - 1)
            col_widths = hm.get("col_widths", default_cw)
            # Normalise widths so they sum to right_w
            total = sum(col_widths)
            col_widths = [w * right_w / total for w in col_widths]

            row_h   = 0.33
            cell_h  = 0.35

            # Default level→colour mapping (overridable via heatmap.level_colors)
            level_cfg = hm.get("level_colors", {})
            level_color_map = {
                "High":   t.color(level_cfg.get("High",   "primary")),
                "Medium": t.color(level_cfg.get("Medium", "orange")),
                "Low":    t.color(level_cfg.get("Low",    "green")),
            }

            # Header row
            cx = right_x
            for ci, hdr in enumerate(headers):
                shp = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(hdr_y), Inches(col_widths[ci]), Inches(row_h)
                )
                shp.fill.solid()
                shp.fill.fore_color.rgb = t.dark
                shp.line.fill.background()
                tf = shp.text_frame
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = hdr
                p.alignment = PP_ALIGN.CENTER
                p.font.size = self._sz(8)
                p.font.bold = True
                p.font.color.rgb = t.white
                cx += col_widths[ci]

            # Data rows
            for ri, row in enumerate(hm_rows):
                ry = hdr_y + row_h + ri * cell_h
                cx = right_x
                values = row.get("values", [])

                # Label cell (first column)
                shp = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(cx), Inches(ry), Inches(col_widths[0]), Inches(cell_h)
                )
                shp.fill.solid()
                shp.fill.fore_color.rgb = t.light_gray
                shp.line.fill.background()
                tf = shp.text_frame
                tf.clear()
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = row.get("label", "")
                p.alignment = PP_ALIGN.CENTER
                p.font.size = self._sz(9.3)
                p.font.color.rgb = t.dark
                cx += col_widths[0]

                # Value cells (colour-coded)
                for vi, val in enumerate(values):
                    ci = vi + 1
                    if ci >= n_cols:
                        break
                    cell_fill = level_color_map.get(val, t.mid_gray)
                    shp = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx), Inches(ry),
                        Inches(col_widths[ci]), Inches(cell_h)
                    )
                    shp.fill.solid()
                    shp.fill.fore_color.rgb = cell_fill
                    shp.line.fill.background()
                    tf = shp.text_frame
                    tf.clear()
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p = tf.paragraphs[0]
                    p.text = val
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = self._sz(9.3)
                    p.font.color.rgb = t.white
                    cx += col_widths[ci]

        # ── bottom row: risk themes (left) + takeaway (right) ─────────────
        hm_rows_count = len(sd.get("heatmap", {}).get("rows", []))
        hdr_y_val     = y_top + 0.32
        bottom_y      = hdr_y_val + 0.33 + hm_rows_count * 0.35 + 0.18

        rt = sd.get("risk_themes", {})
        if rt:
            rt_h   = rt.get("height", 1.65)
            rt_fill = t.color(rt.get("color", "orange"))
            self._card(
                slide, left_x, bottom_y, left_w, rt_h,
                rt.get("title", "Primary risk themes"),
                rt.get("items", []),
                fill=t.white, accent=rt_fill,
                title_size=rt.get("title_size", 14.5),
                bullet_size=rt.get("bullet_size", 10.3),
            )

        ta = sd.get("takeaway", {})
        if ta:
            ta_h = ta.get("height", 1.25)
            self._card(
                slide, right_x, bottom_y, right_w, ta_h,
                ta.get("title", "Executive takeaway"),
                [ta.get("text", "")],
                fill=t.light_gray, accent=t.primary,
                title_size=ta.get("title_size", 14.5),
                bullet_size=ta.get("bullet_size", 11.2),
            )

    def _render_dependency(self, sd: dict):
        """
        Three-column dependency/flow diagram.
        Left → Middle → Right with Unicode arrow separators.
        Implements the decision dependency pattern from
        decisiondependencydiagramdeck.py.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        nodes = sd.get("nodes", [])
        if not nodes:
            return

        n_rows = len(nodes)
        header_y = 1.38
        header_h = 0.40

        # Column layout — room for arrows between cols
        lx, lw = 0.55, 3.0
        ax1 = lx + lw + 0.02
        mx, mw = ax1 + 0.52, 4.4
        ax2 = mx + mw + 0.02
        rx, rw = ax2 + 0.52, 13.333 - (ax2 + 0.52) - 0.55

        # Header row
        self._badge(
            slide, lx, header_y, lw, header_h,
            sd.get("col1_label", "Input / Decision"), t.dark, size=9
        )
        self._badge(
            slide, mx, header_y, mw, header_h,
            sd.get("col2_label", "Work Item / Action"), t.mid_gray, size=9
        )
        self._badge(
            slide, rx, header_y, rw, header_h,
            sd.get("col3_label", "Outcome"), t.dark, size=9
        )

        row_y = header_y + header_h + 0.12
        available_h = 6.85 - row_y
        row_h = min(available_h / n_rows, 0.78) - 0.1

        for node in nodes:
            fill = t.color(node.get("color", "primary"))
            mid_y = row_y + row_h / 2

            self._badge(slide, lx, row_y, lw, row_h, node.get("left", ""), fill, size=9)
            self._textbox(
                slide, ax1, mid_y - 0.16, 0.52, 0.38, "→",
                size=20, color=t.mid_gray, bold=True, align=PP_ALIGN.CENTER
            )
            self._badge(
                slide, mx, row_y, mw, row_h, node.get("middle", ""),
                t.light_gray, font_color=t.dark, size=9
            )
            self._textbox(
                slide, ax2, mid_y - 0.16, 0.52, 0.38, "→",
                size=20, color=t.mid_gray, bold=True, align=PP_ALIGN.CENTER
            )
            self._badge(slide, rx, row_y, rw, row_h, node.get("right", ""), fill, size=9)

            row_y += row_h + 0.12

    def _render_timeline(self, sd: dict):
        """
        Vertical phase columns: Now / Next / Later (or any custom phases).
        Each phase gets a colour header band + optional period label +
        individual item mini-cards.
        Implements the pattern from recommendedfuturetemplatescript.py.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        phases = sd.get("phases", [])
        if not phases:
            return

        n = len(phases)
        margin, gap = 0.55, 0.22
        pw = (13.333 - 2 * margin - gap * (n - 1)) / n
        py = 1.38
        default_accents = ["primary", "blue", "green", "teal", "orange", "navy"]

        for i, phase in enumerate(phases):
            x = margin + i * (pw + gap)
            accent = t.color(phase.get("accent", default_accents[i % len(default_accents)]))

            # Phase header band
            self._badge(slide, x, py, pw, 0.52, phase.get("title", f"Phase {i+1}"), accent, size=13)

            item_y = py + 0.62
            if phase.get("period"):
                self._textbox(
                    slide, x + 0.05, item_y, pw - 0.1, 0.28,
                    phase["period"], size=9, color=t.mid_gray, align=PP_ALIGN.CENTER
                )
                item_y += 0.32

            for bullet in phase.get("bullets", []):
                if item_y > 6.8:
                    break
                self._card(
                    slide, x + 0.05, item_y, pw - 0.1, 0.48,
                    bullet, fill=t.light_gray, accent=accent,
                    title_size=11, bullet_size=10,
                )
                item_y += 0.55

    def _render_quote(self, sd: dict):
        """
        Large quote / statement slide.
        Decorative opening-quotation-mark, bold text, optional attribution.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)

        qm = slide.shapes.add_textbox(Inches(0.35), Inches(0.3), Inches(2), Inches(2))
        p = qm.text_frame.paragraphs[0]
        p.text = "\u201c"
        p.font.size = self._sz(112)
        p.font.bold = True
        p.font.color.rgb = t.light_gray

        tb = slide.shapes.add_textbox(
            Inches(1.1), Inches(1.65), Inches(11.2), Inches(3.3)
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p2 = tf.paragraphs[0]
        p2.text = sd.get("quote", "")
        p2.font.size = self._sz(30)
        p2.font.bold = True
        p2.font.color.rgb = t.dark

        if sd.get("attribution"):
            at = slide.shapes.add_textbox(
                Inches(1.1), Inches(5.1), Inches(11.2), Inches(0.5)
            )
            p3 = at.text_frame.paragraphs[0]
            p3.text = f"\u2014 {sd['attribution']}"
            p3.font.size = self._sz(16)
            p3.font.color.rgb = t.mid_gray

    def _render_metrics(self, sd: dict):
        """
        Big-number KPI boxes.
        Each metric has value (large number), label, and optional sub-text.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        metrics = sd.get("metrics", [])
        n = len(metrics)
        if n == 0:
            return

        margin, gap = 0.55, 0.3
        bw = (13.333 - 2 * margin - gap * (n - 1)) / n
        bh, by = 3.8, 1.85
        default_colors = ["primary", "blue", "teal", "green", "orange", "navy"]

        for i, m in enumerate(metrics):
            x = margin + i * (bw + gap)
            fill = t.color(m.get("color", default_colors[i % len(default_colors)]))

            shp = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(by), Inches(bw), Inches(bh)
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
            shp.line.fill.background()

            tf = shp.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.text = m.get("value", "—")
            p.font.size = self._sz(54)
            p.font.bold = True
            p.font.color.rgb = t.white
            p.alignment = PP_ALIGN.CENTER

            lbl = tf.add_paragraph()
            lbl.text = m.get("label", "")
            lbl.font.size = self._sz(16)
            lbl.font.color.rgb = t.white
            lbl.alignment = PP_ALIGN.CENTER

            if m.get("sub"):
                sub = tf.add_paragraph()
                sub.text = m["sub"]
                sub.font.size = self._sz(12)
                sub.font.color.rgb = t.white
                sub.alignment = PP_ALIGN.CENTER

        if sd.get("note"):
            self._textbox(
                slide, 0.55, by + bh + 0.2, 12.2, 0.4, sd["note"],
                size=10, color=t.mid_gray, align=PP_ALIGN.CENTER
            )

    def _render_divider(self, sd: dict):
        """
        Section-break divider slide.
        White background, primary-colour top + bottom bars,
        primary-colour title text, dark subtitle text.
        """
        t = self.t
        slide = self._blank()

        # White background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        bg.fill.solid()
        bg.fill.fore_color.rgb = t.white
        bg.line.fill.background()

        # Top primary bar
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.35)
        )
        top.fill.solid()
        top.fill.fore_color.rgb = t.primary
        top.line.fill.background()

        # Bottom primary bar
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25)
        )
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = t.primary
        bottom.line.fill.background()

        if sd.get("section_number"):
            sn = slide.shapes.add_textbox(
                Inches(0.65), Inches(2.0), Inches(5), Inches(0.4)
            )
            p3 = sn.text_frame.paragraphs[0]
            p3.text = f"SECTION {sd['section_number']}"
            p3.font.size = self._sz(13)
            p3.font.bold = True
            p3.font.color.rgb = t.mid_gray

        tb = slide.shapes.add_textbox(
            Inches(0.65), Inches(2.5), Inches(12.0), Inches(1.8)
        )
        tf = tb.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sd.get("title", "")
        p.font.size = self._sz(46)
        p.font.bold = True
        p.font.color.rgb = t.primary

        if sd.get("subtitle"):
            st = slide.shapes.add_textbox(
                Inches(0.65), Inches(4.45), Inches(12.0), Inches(0.55)
            )
            p2 = st.text_frame.paragraphs[0]
            p2.text = sd["subtitle"]
            p2.font.size = self._sz(20)
            p2.font.color.rgb = t.dark

    def _render_table(self, sd: dict):
        """
        Simple styled data table with a coloured header row and
        alternating light-gray row shading.
        """
        t = self.t
        slide = self._blank()
        self._bg(slide)
        self._heading(slide, sd.get("title", ""), sd.get("subtitle"))

        headers = sd.get("headers", [])
        rows = sd.get("rows", [])
        if not headers or not rows:
            return

        n_cols = len(headers)
        n_rows = len(rows)
        table_h = min(5.5, n_rows * 0.50 + 0.55)

        tbl = slide.shapes.add_table(
            n_rows + 1, n_cols,
            Inches(0.5), Inches(1.48),
            Inches(12.35), Inches(table_h)
        ).table

        # Header row
        for ci, hdr in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = str(hdr)
            cell.fill.solid()
            cell.fill.fore_color.rgb = t.primary
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = self._sz(13)
            p.font.color.rgb = t.white
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = t.light_gray
                p = cell.text_frame.paragraphs[0]
                p.font.size = self._sz(12)
                p.font.color.rgb = t.dark

    # -----------------------------------------------------------------------
    # SLIDE TYPE DISPATCH
    # -----------------------------------------------------------------------

    _RENDERERS: dict[str, str] = {
        # Title / cover
        "title":            "_render_title",
        "cover":            "_render_title",
        # Bullet lists
        "bullets":          "_render_bullets",
        "bullet":           "_render_bullets",
        "list":             "_render_bullets",
        # Card columns
        "cards":            "_render_cards",
        "columns":          "_render_cards",
        "three_column":     "_render_cards",
        "two_column":       "_render_cards",
        # Decision + risk
        "decision_risk":    "_render_decision_risk",
        "decisions":        "_render_decision_risk",
        "risk":             "_render_decision_risk",
        "rag":              "_render_decision_risk",
        # Card grid (multi-row grid or two-section panel)
        "card_grid":        "_render_card_grid",
        "grid":             "_render_card_grid",
        "multi_row":        "_render_card_grid",
        "two_section":      "_render_card_grid",
        # Decisions + coloured risk heatmap
        "decisions_heatmap":  "_render_decisions_heatmap",
        "risk_heatmap":       "_render_decisions_heatmap",
        "heatmap":            "_render_decisions_heatmap",
        # Dependency / flow
        "dependency":       "_render_dependency",
        "flow":             "_render_dependency",
        "pipeline":         "_render_dependency",
        # Timeline / roadmap / phases
        "timeline":         "_render_timeline",
        "roadmap":          "_render_timeline",
        "phases":           "_render_timeline",
        "now_next_later":   "_render_timeline",
        # Quote / statement
        "quote":            "_render_quote",
        "statement":        "_render_quote",
        "callout":          "_render_quote",
        # Metrics / KPIs
        "metrics":          "_render_metrics",
        "kpi":              "_render_metrics",
        "numbers":          "_render_metrics",
        # Section dividers
        "divider":          "_render_divider",
        "section":          "_render_divider",
        "chapter":          "_render_divider",
        # Data tables
        "table":            "_render_table",
        "data":             "_render_table",
    }

    def render(self, slide_def: dict):
        """Dispatch to the correct renderer based on the 'type' key."""
        raw_type = slide_def.get("type", "bullets")
        stype = str(raw_type).lower().replace("-", "_").replace(" ", "_")
        method_name = self._RENDERERS.get(stype)
        if method_name:
            getattr(self, method_name)(slide_def)
        else:
            print(f"    [warn] Unknown slide type '{raw_type}' — falling back to bullets")
            self._render_bullets(slide_def)

    def save(self, path: str):
        self.prs.save(path)


# ===========================================================================
# PARSERS
# ===========================================================================

def _parse_yaml_file(path: Path) -> dict:
    """Load a YAML deck definition."""
    with open(path, encoding="utf-8") as f:
        data = yaml.safe_load(f) or {}
    return data


def _parse_markdown_file(path: Path) -> dict:
    """
    Parse a structured Markdown file into a deck definition dict.

    Format overview:
    ─────────────────────────────────────────────────
    ---
    title: My Deck
    footer: Company | Dept | 2025
    output: my_deck.pptx
    theme: corporate_red
    ---

    # Slide Title [type:cards]
    > Optional subtitle text

    ### Column One [accent:primary]
    - Bullet one
    - Bullet two

    ### Column Two [accent:blue]
    - Bullet one

    ---

    # Another Slide [type:bullets]
    - Bullet one
      - Sub bullet
    - Bullet two
    ─────────────────────────────────────────────────
    """
    with open(path, encoding="utf-8") as f:
        content = f.read()

    deck: dict = {"meta": {}, "slides": []}

    # YAML front matter
    fm = re.match(r"^---\s*\n(.*?)\n---\s*\n", content, re.DOTALL)
    if fm:
        deck["meta"] = yaml.safe_load(fm.group(1)) or {}
        content = content[fm.end():]

    for raw_block in re.split(r"\n---\n", content):
        raw_block = raw_block.strip()
        if raw_block:
            sd = _parse_md_block(raw_block)
            if sd:
                deck["slides"].append(sd)

    return deck


def _parse_md_block(raw: str) -> dict:
    """Parse a single Markdown slide block."""
    lines = raw.split("\n")
    sd: dict = {}

    # First line must be a heading: # Title [key:val] [key:val]
    m = re.match(r"^#{1,3}\s+(.*?)(\s+\[([^\]]+)\])?\s*$", lines[0])
    if not m:
        return {}

    sd["title"] = m.group(1).strip()
    sd["type"] = "bullets"

    if m.group(3):
        for directive in m.group(3).split():
            if ":" in directive:
                k, v = directive.split(":", 1)
                sd[k] = v

    rest = lines[1:]

    stype = sd["type"]

    # Optional subtitle (blockquote immediately after heading).
    # Skip for quote-type slides so the blockquote stays available for the
    # quote-specific handler below (otherwise it disappears into 'subtitle').
    if rest and rest[0].strip().startswith(">") and stype not in ("quote", "statement", "callout"):
        sd["subtitle"] = rest[0].strip().lstrip("> ").strip()
        rest = rest[1:]

    if stype in ("bullets", "bullet", "list"):
        bullets = []
        for line in rest:
            bm = re.match(r"^(\s*)[-*+]\s+(.+)", line)
            if bm:
                depth = len(bm.group(1)) // 2
                prefix = "  " * depth
                bullets.append(prefix + bm.group(2).rstrip())
        sd["bullets"] = bullets

    elif stype in ("cards", "columns", "three_column", "two_column"):
        cols = []
        cur: dict | None = None
        for line in rest:
            h3 = re.match(r"^###\s+(.*?)(\s+\[([^\]]+)\])?\s*$", line)
            if h3:
                if cur is not None:
                    cols.append(cur)
                cur = {"title": h3.group(1).strip(), "bullets": []}
                if h3.group(3):
                    for d in h3.group(3).split():
                        if ":" in d:
                            k, v = d.split(":", 1)
                            cur[k] = v
            elif cur is not None:
                bm = re.match(r"^\s*[-*+]\s+(.+)", line)
                if bm:
                    cur["bullets"].append(bm.group(1).rstrip())
                elif line.strip().startswith(">"):
                    sd["takeaway"] = line.strip().lstrip("> ").strip()
        if cur is not None:
            cols.append(cur)
        sd["columns"] = cols

    elif stype in ("quote", "statement", "callout"):
        for line in rest:
            ls = line.strip()
            if ls.startswith(">"):
                sd["quote"] = ls.lstrip("> ").strip()
            elif ls.startswith("—") or ls.startswith("--"):
                sd["attribution"] = ls.lstrip("—- ").strip()

    # Other types (title, divider, section, metrics, dependency, etc.)
    # require YAML for full fidelity; basic title is already captured above.

    return sd


# ===========================================================================
# BUILD
# ===========================================================================

def build(deck_def: dict, *, output_override: str = None, theme_override: str = None,
          base_dir: Path = None, input_path: Path = None, text_offset: int = 0) -> str:
    """Render all slides and save the presentation. Returns the output path."""
    meta = deck_def.get("meta") or {}
    theme_cfg = theme_override or meta.get("theme", "Corp_red")
    theme = Theme(theme_cfg)
    footer = meta.get("footer", "")

    # Resolve logo path.  Priority:
    #   1. meta.logo in the YAML (resolved relative to the input file)
    #   2. Default bundled logo for the Corp_red theme
    _DEFAULT_LOGO = Path(__file__).parent / "assets" / "logo.png"
    logo_path = None
    raw_logo = meta.get("logo")
    if raw_logo:
        resolve_base = base_dir or Path.cwd()
        logo_path = str((resolve_base / raw_logo).resolve())
    elif isinstance(theme_cfg, str) and theme_cfg == "Corp_red" and _DEFAULT_LOGO.exists():
        logo_path = str(_DEFAULT_LOGO)

    renderer = DeckRenderer(theme, footer, logo_path=logo_path, text_offset=text_offset)

    slides = deck_def.get("slides", [])
    total = len(slides)

    for i, sd in enumerate(slides):
        label = sd.get("title", "(untitled)")[:60]
        stype = sd.get("type", "bullets")
        print(f"  [{i + 1:2d}/{total}]  {stype:<20}  {label}")
        renderer.render(sd)

    default_name = (input_path.stem + ".pptx") if input_path else "output.pptx"
    out_path = output_override or meta.get("output", default_name)
    renderer.save(out_path)
    return out_path


# ===========================================================================
# TEMPLATE GENERATORS
# ===========================================================================

_YAML_TEMPLATE = """\
# ============================================================
# DeckBuilder — Starter YAML Template
# All twelve slide types are shown below.
# Run: python deckbuilder.py this_file.yaml
# ============================================================

meta:
  title: "My Presentation"
  footer: "Organisation | Department | 2025"
  output: "output.pptx"
  theme: Corp_red    # Options: Corp_red, corporate_red, corporate_blue,
                          #          corporate_green, modern_purple, minimal_slate,
                          #          executive_dark
  # logo: "assets/logo.png"   # Optional: path to logo image (relative to this file)

slides:

  # ── 1. TITLE ─────────────────────────────────────────────
  # Full-bleed dark cover slide.
  - type: title
    title: "Presentation Title"
    subtitle: "Optional subtitle line"               # optional
    context: "Author · Organisation · Month Year"   # optional small attribution line

  # ── 2. DIVIDER ───────────────────────────────────────────
  # Dark section-break slide.
  - type: divider
    section_number: 1           # optional — renders "SECTION 1" above the title
    title: "Section Name"
    subtitle: "Brief description of this section"   # optional
    bg_color: navy              # optional — palette key or #hex (default: navy)

  # ── 3. BULLETS ───────────────────────────────────────────
  # Standard content slide with a bullet list (up to 3 indent levels).
  - type: bullets
    title: "Slide Title"
    subtitle: "Optional subtitle"                    # optional
    intro: "Optional bold introductory sentence."   # optional
    bullets:
      - "Top-level bullet"
      - "Another top-level bullet"
      - "  Sub-bullet (2 leading spaces = level 1)"
      - "    Deeper sub-bullet (4 spaces = level 2)"
      # Alternative explicit syntax:
      # - text: "Explicit item"
      #   level: 0

  # ── 4. CARDS ─────────────────────────────────────────────
  # Two-to-four column card layout with coloured accent bars.
  - type: cards
    title: "Strategic Pillars"
    subtitle: "Optional subtitle"
    intro: "Optional bold introductory sentence above the cards."
    # card_height: 4.8   # optional — card height in inches (default: auto)
    columns:
      - title: "Column One"
        accent: primary    # palette key or #hex
        fill: white        # card background (default: white)
        bullets:
          - "First point"
          - "Second point"
      - title: "Column Two"
        accent: blue
        bullets:
          - "First point"
          - "Second point"
      - title: "Column Three"
        accent: green
        bullets:
          - "First point"
          - "Second point"
    takeaway: "Optional full-width summary card below the columns."   # optional

  # ── 5. DECISION_RISK ─────────────────────────────────────
  # Decision boxes + RAG (Red/Amber/Green) risk grid.
  - type: decision_risk
    title: "Key Decisions & Risk Assessment"
    subtitle: "Owners, timelines, and leadership alignment needed"
    decisions:
      - title: "Decision A"
        body: "Owner: Team\nTimeline: Q3\nContext or options here"
        color: primary
      - title: "Decision B"
        body: "Owner: Team\nTimeline: Q4\nContext or options here"
        color: blue
    risks:
      high_color: primary      # optional — default: primary
      high:
        - "High-risk item one"
        - "High-risk item two"
      medium_color: orange     # optional — default: orange
      medium:
        - "Medium-risk item one"
      low_color: green         # optional — default: green
      low:
        - "Low-risk item one"

  # ── 6. CARD_GRID ─────────────────────────────────────────
  # Multi-row card grid or two-section side-by-side panel layout.
  # Mode 1 — simple grid (rows key):
  - type: card_grid
    title: "Delivery Summary"
    subtitle: "Multi-row card layout"
    intro: "Optional bold introductory sentence above the cards."
    # header_badge:            # optional centred badge above all cards
    #   title: "Key Outcome"
    #   color: primary
    rows:
      - columns:
          - title: "Row 1 Card 1"
            accent: primary
            fill: white
            bullets:
              - "First point"
              - "Second point"
          - title: "Row 1 Card 2"
            accent: blue
            bullets:
              - "First point"
              - "Second point"
          - title: "Row 1 Card 3"
            accent: green
            bullets:
              - "First point"
              - "Second point"
      - columns:
          - title: "Row 2 Card 1"
            accent: teal
            bullets:
              - "First point"
              - "Second point"
          - title: "Row 2 Card 2"
            accent: orange
            bullets:
              - "First point"
              - "Second point"
          - title: "Row 2 Card 3"
            accent: navy
            bullets:
              - "First point"
              - "Second point"
    # badges:                  # optional outcome pills below cards
    #   - "Outcome A"
    #   - "Outcome B"
    takeaway: "Optional full-width summary card at the bottom."   # optional
    # note: "Optional footer note text"                           # optional

  # ── 7. DECISIONS_HEATMAP ─────────────────────────────────
  # Decisions panel + colour-coded risk heatmap table (side-by-side).
  - type: decisions_heatmap
    title: "Decisions & Risk Heatmap"
    subtitle: "Side-by-side decisions and colour-coded risk matrix"
    decisions:
      title: "Strategic decisions required"
      color: primary        # accent bar colour
      items:
        - "Approve Q3 budget allocation"
        - "Confirm vendor selection"
        - "Finalise headcount plan"
    risk_themes:             # optional — second box below decisions
      title: "Primary risk themes"
      color: orange
      items:
        - "Resource contention across workstreams"
        - "Regulatory timeline uncertainty"
    heatmap:
      title: "Risk Heatmap Summary"
      headers:
        - "Risk Area"
        - "Financial"
        - "Schedule"
        - "Resource"
        - "Technical"
      rows:
        - label: "Workstream A"
          values: ["High", "Medium", "Low", "Medium"]
        - label: "Workstream B"
          values: ["Medium", "High", "Medium", "Low"]
        - label: "Workstream C"
          values: ["Low", "Low", "High", "Medium"]
      # level_colors:         # optional — override default High/Medium/Low colours
      #   High: primary
      #   Medium: orange
      #   Low: green
    takeaway:                # optional summary card below heatmap
      title: "Executive takeaway"
      text: "Key risk mitigation actions and timeline implications."

  # ── 8. DEPENDENCY ────────────────────────────────────────
  # Three-column left→middle→right flow / dependency diagram.
  - type: dependency
    title: "Decision Dependency Sequence"
    subtitle: "How decisions chain into outcomes"
    col1_label: "Input / Decision"     # optional column header labels
    col2_label: "Work Item / Action"
    col3_label: "Outcome"
    nodes:
      - left: "Input A\\nQ3 | Team Name"
        middle: "Action or work item description"
        right: "Expected outcome\\nor result"
        color: teal
      - left: "Input B\\nQ4 | Team Name"
        middle: "Second action description"
        right: "Second outcome"
        color: primary

  # ── 9. TIMELINE ──────────────────────────────────────────
  # Phase-based roadmap columns (Now / Next / Later or quarterly).
  - type: timeline
    title: "Roadmap 2025–2027"
    subtitle: "Three-horizon view"
    phases:
      - title: "Now  (H2 2025)"
        accent: primary
        period: "Jul – Dec 2025"    # optional small label under header
        bullets:
          - "Deliverable one"
          - "Deliverable two"
      - title: "Next  (2026)"
        accent: blue
        period: "Jan – Dec 2026"
        bullets:
          - "Deliverable one"
          - "Deliverable two"
      - title: "Later  (2027+)"
        accent: teal
        bullets:
          - "Deliverable one"
          - "Deliverable two"

  # ── 10. METRICS ──────────────────────────────────────────
  # Large KPI / big-number boxes.
  - type: metrics
    title: "Key Performance Indicators"
    subtitle: "Q2 2025 · Platform-level outcomes"
    metrics:
      - value: "94%"
        label: "Metric Label"
        sub: "+12pp vs prior period"   # optional sub-text
        color: primary
      - value: "47"
        label: "Another Metric"
        color: blue
      - value: "3.2×"
        label: "Third Metric"
        color: teal
    note: "Source: Dashboard · Data as at 30 June 2025"   # optional footer note

  # ── 11. QUOTE ────────────────────────────────────────────
  # Large bold statement slide with decorative quotation mark.
  - type: quote
    quote: "Replace this with your key insight or executive statement."
    attribution: "Author Name, Role · Year"   # optional

  # ── 12. TABLE ────────────────────────────────────────────
  # Styled data table with coloured header row.
  - type: table
    title: "Service Inventory"
    subtitle: "Key items by status"
    headers:
      - "Name"
      - "Owner"
      - "Status"
      - "Target"
      - "Date"
    rows:
      - ["Row 1 Col 1", "Owner A", "In Progress",    "Target A", "Q3 2025"]
      - ["Row 2 Col 1", "Owner B", "Decision Needed","Target B", "Q4 2025"]
      - ["Row 3 Col 1", "Owner C", "Complete",        "Target C", "Q1 2026"]
"""

_MD_TEMPLATE = """\
---
title: "My Presentation"
footer: "Organisation | Department | 2025"
output: "output.pptx"
theme: Corp_red
# Available themes: Corp_red, corporate_red, corporate_blue,
#                   corporate_green, modern_purple, minimal_slate, executive_dark
---

# Presentation Title [type:title]
> Optional subtitle line
— Author · Organisation · Month Year

---

# Section One [type:divider]

---

# Slide Title [type:bullets]
> Optional subtitle

- Top-level bullet
- Another top-level bullet
  - Sub-bullet (2 leading spaces = level 1)
    - Deeper sub-bullet (4 spaces = level 2)

---

# Three Pillars [type:cards]
> Optional subtitle for cards slide

### Column One [accent:primary]
- First point
- Second point

### Column Two [accent:blue]
- First point
- Second point

### Column Three [accent:green]
- First point
- Second point

> Optional full-width takeaway card below all columns.

---

# Section Two [type:divider]

---

# Roadmap [type:bullets]
> Now · Next · Later execution priorities
> Note: for full timeline/roadmap slides use YAML (type:timeline)

- NOW (H2 2025): Deliverable one, deliverable two
  - Supporting item
- NEXT (2026): Deliverable one, deliverable two
  - Supporting item
- LATER (2027+): Future deliverable
  - Supporting item

---

# Key Insight [type:quote]
> Replace this with your executive statement or key insight.
— Author Name, Role · Year

---

# Note on Markdown limitations
# The following slide types require YAML for full fidelity:
#   dependency, decision_risk, card_grid, decisions_heatmap, metrics/kpi,
#   timeline/roadmap, table
# Use 'deckbuilder --init-yaml starter.yaml' to get the full YAML template.
"""


def _create_yaml_template(dest: Path) -> None:
    """Write a starter YAML template covering all twelve slide types."""
    if dest.exists():
        sys.exit(f"Error: '{dest}' already exists — refusing to overwrite.")
    dest.write_text(_YAML_TEMPLATE, encoding="utf-8")
    print(f"Created YAML template → {dest}")
    print("Edit the file, then run:  deckbuilder " + str(dest))


def _create_md_template(dest: Path) -> None:
    """Write a starter Markdown template covering Markdown-compatible slide types."""
    if dest.exists():
        sys.exit(f"Error: '{dest}' already exists — refusing to overwrite.")
    dest.write_text(_MD_TEMPLATE, encoding="utf-8")
    print(f"Created Markdown template → {dest}")
    print("Edit the file, then run:  deckbuilder " + str(dest))


def _print_readme() -> None:
    """Locate and print the README.md bundled alongside this script.

    Renders headings, fenced code blocks, horizontal rules, and inline
    bold/code markers with ANSI colours so the content is readable in any
    terminal without requiring third-party tools.
    """
    import shutil
    import re

    readme = Path(__file__).with_name("README.md")
    if not readme.exists():
        sys.exit("Error: README.md not found alongside deckbuilder.py")

    text = readme.read_text(encoding="utf-8")

    # --- Structured plain-text renderer -----------------------------------
    term_width = shutil.get_terminal_size(fallback=(100, 40)).columns
    col = min(term_width, 100)

    RESET  = "\033[0m"
    BOLD   = "\033[1m"
    DIM    = "\033[2m"
    CYAN   = "\033[36m"
    YELLOW = "\033[33m"
    GREEN  = "\033[32m"

    def _strip_inline(s: str) -> str:
        """Remove bold/italic markdown markers for plain rendering."""
        s = re.sub(r"\*\*(.+?)\*\*", rf"{BOLD}\1{RESET}", s)
        s = re.sub(r"\*(.+?)\*",     rf"\1",               s)
        s = re.sub(r"`(.+?)`",       rf"{GREEN}\1{RESET}",  s)
        return s

    print()
    in_code = False
    code_lang = ""

    for line in text.splitlines():
        # Fenced code blocks
        fence = re.match(r"^```(\w*)", line)
        if fence:
            if not in_code:
                in_code = True
                code_lang = fence.group(1)
                label = f" {code_lang} " if code_lang else ""
                print(f"{DIM}{'─'*4}{label}{'─'*(col - 6 - len(label))}{RESET}")
            else:
                in_code = False
                print(f"{DIM}{'─' * col}{RESET}")
            continue

        if in_code:
            print(f"  {GREEN}{line}{RESET}")
            continue

        # Horizontal rules
        if re.match(r"^[-*_]{3,}\s*$", line):
            print(f"{DIM}{'─' * col}{RESET}")
            continue

        # ATX headings
        m = re.match(r"^(#{1,6})\s+(.*)", line)
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            if level == 1:
                bar = "═" * col
                print(f"\n{CYAN}{BOLD}{bar}{RESET}")
                print(f"{CYAN}{BOLD}  {title}{RESET}")
                print(f"{CYAN}{BOLD}{bar}{RESET}")
            elif level == 2:
                print(f"\n{YELLOW}{BOLD}  {title}{RESET}")
                print(f"{YELLOW}  {'─' * (col - 2)}{RESET}")
            elif level == 3:
                print(f"\n{BOLD}    {title}{RESET}")
            else:
                print(f"\n{BOLD}      {'#'*(level-3)} {title}{RESET}")
            continue

        # Blank lines
        if not line.strip():
            print()
            continue

        # Bullet / numbered list items — indent visually
        m_bullet = re.match(r"^(\s*)([-*+]|\d+\.)\s+(.*)", line)
        if m_bullet:
            indent = m_bullet.group(1)
            marker = "•" if re.match(r"[-*+]", m_bullet.group(2)) else m_bullet.group(2)
            body   = _strip_inline(m_bullet.group(3))
            print(f"  {indent}{BOLD}{marker}{RESET}  {body}")
            continue

        # Normal paragraph text
        print(f"  {_strip_inline(line)}")

    print()


# ===========================================================================
# CLI
# ===========================================================================

def main():
    ap = argparse.ArgumentParser(
        prog="deckbuilder",
        description="Convert YAML or Markdown to polished PowerPoint presentations.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""
            Examples:
              python deckbuilder.py deck.yaml
              python deckbuilder.py deck.md
              python deckbuilder.py deck.yaml -o quarterly_review.pptx
              python deckbuilder.py deck.yaml --theme corporate_blue
              python deckbuilder.py --list-themes
              python deckbuilder.py --init-yaml starter.yaml
              python deckbuilder.py --init-md   starter.md
              python deckbuilder.py --help-full

            Slide types:
              title, divider, bullets, cards, decision_risk,
              dependency, timeline, quote, metrics, table

            Input formats:
              .yaml / .yml   Full feature support (recommended)
              .md / .markdown  Good for bullets + cards; use YAML for advanced types
        """),
    )
    ap.add_argument("input", nargs="?", help="Input YAML or Markdown file")
    ap.add_argument("-o", "--output", metavar="FILE", help="Output .pptx path")
    ap.add_argument("--theme", metavar="NAME", help="Override theme (see --list-themes)")
    ap.add_argument(
        "--list-themes", action="store_true",
        help="Print available built-in themes and exit"
    )
    ap.add_argument(
        "--init-yaml", metavar="FILE", nargs="?", const="starter.yaml",
        help="Create a starter YAML template (default: starter.yaml) and exit"
    )
    ap.add_argument(
        "--init-md", metavar="FILE", nargs="?", const="starter.md",
        help="Create a starter Markdown template (default: starter.md) and exit"
    )
    ap.add_argument(
        "--help-full", action="store_true",
        help="Display the full README.md documentation and exit"
    )
    ap.add_argument(
        "--text-size-offset", metavar="N", type=int, default=0,
        help="Adjust all text sizes by N points (positive = larger, negative = smaller)"
    )
    args = ap.parse_args()

    if args.help_full:
        _print_readme()
        return

    if args.list_themes:
        print("Built-in themes:")
        for name, palette in BUILT_IN_THEMES.items():
            primary = palette.get("primary", "?")
            print(f"  {name:<22}  primary={primary}")
        return

    if args.init_yaml is not None:
        p = Path(args.init_yaml)
        if p.suffix.lower() not in (".yaml", ".yml"):
            p = p.with_suffix(".yaml")
        _create_yaml_template(p)
        return

    if args.init_md is not None:
        p = Path(args.init_md)
        if p.suffix.lower() not in (".md", ".markdown"):
            p = p.with_suffix(".md")
        _create_md_template(p)
        return

    if not args.input:
        ap.print_help()
        sys.exit(1)

    src = Path(args.input)
    if not src.exists():
        sys.exit(f"Error: file not found: {src}")

    sep = "─" * 54
    print(f"\nDeckBuilder\n{sep}\nSource  : {src}")

    ext = src.suffix.lower()
    if ext in (".yaml", ".yml"):
        deck_def = _parse_yaml_file(src)
    elif ext in (".md", ".markdown"):
        deck_def = _parse_markdown_file(src)
    else:
        sys.exit(f"Error: unsupported format '{ext}' — use .yaml, .yml, or .md")

    theme_name = args.theme or (deck_def.get("meta") or {}).get("theme", "Corp_red")
    print(f"Theme   : {theme_name}")
    print(f"{sep}\n")

    out = build(deck_def, output_override=args.output, theme_override=args.theme,
                base_dir=src.parent, input_path=src,
                text_offset=args.text_size_offset)
    print(f"\n{sep}\nSaved   → {out}\n")


if __name__ == "__main__":
    main()
